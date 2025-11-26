"""Extract jaguar identification data from PowerPoint presentations.

This script processes PowerPoint files containing jaguar sighting information and images,
extracting metadata and saving images to create a dataset compatible with the main labels format.

The script:
- Extracts text fields (NOME, NUMBER ID, LOCALIZAÇÃO, SITE, SEX, DATAS, OBS) from each slide
- Saves all images from slides to the output directory
- Creates one row per image in the output CSV
- Outputs a CSV file with the same format as clean_labels.py for easy dataset merging

Output CSV columns:
- CAMERA TRAP SITE, LATITUDE, LONGITUDE, CAMERA ID, CAM, JAGUAR ID, LOCATION, CAMERA MODEL,
  DATE, TIME, TEMP C, Files Name, DATETIME, ORIGINAL FILE NAME, SIGHTING ID, RAW DATA PATH,
  FILE PATH, FILE TYPE, FILE EXTENSION, FILE NAME, ORIGINAL CAM, ORIGINAL SITE

Run as a module:
    python -m src.jaguar_reidentification.data_preprocessing.pptx_extract \
        --input_pptx="data/raw/17_11_2025/CAMERA TRAP ID GUIDE UPDATED by Oscar2025.pptx" \
        --output_csv=data/intermediate/v1/pptx_extracted_labels.csv \
        --output_dir=data/intermediate/v1/files --generate_report
"""

import argparse
import io
import json
import logging
import re
import unicodedata
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.slide import Slide

from src.jaguar_reidentification.utils.utils import json_safe


def write_report(report: dict, report_path: Path) -> None:
    logging.info("Writing processing report to %s", report_path)
    report = json_safe(report)
    with open(report_path, "w") as f:
        json.dump(report, f, indent=4)


def write_df(df: pd.DataFrame, output_csv: Path) -> dict:
    logging.info("Writing extracted labels to %s", output_csv)
    output_csv.parent.mkdir(parents=True, exist_ok=True)
    df.to_csv(output_csv, index=False)

    report = {
        "row_count": len(df),
        "columns": df.columns.tolist(),
        "unique_jaguar_ids": df["JAGUAR ID"].nunique(),
        "file_types": df["FILE TYPE"].value_counts().to_dict(),
        "file_extensions": df["FILE EXTENSION"].value_counts().to_dict(),
    }
    logging.info("Saved %d extracted rows", report["row_count"])
    return report


def extract_dates(text: str) -> list[str]:
    """Extract date information from slide text using regex patterns."""
    date_patterns = [
        r"(\d{1,2}[\/\.-]\d{1,2}[\/\.-]\d{2,4})",  # DD/MM/YYYY or DD-MM-YYYY or DD.MM.YYYY
        r"(\d{4}[\/\.-]\d{1,2}[\/\.-]\d{1,2})",  # YYYY/MM/DD or YYYY-MM-DD or YYYY.MM.DD
        r"(\d{1,2}\s+[A-Za-z]{3,9}\s+\d{2,4})",  # DD Month YYYY
        r"([A-Za-z]{3,9}\s+\d{1,2},\s+\d{4})",  # Month DD, YYYY
    ]

    dates_found = []
    for pattern in date_patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            dates_found.append(match.strip())

    return dates_found


def extract_fields(text: str) -> dict:
    """Extract jaguar metadata fields from slide text using regex patterns."""
    patterns = {
        "nome": r"NOME\s*:\s*(.*?)(?=LOCALIZAÇÃO|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "name": r"NAME\s*:\s*(.*?)(?=LOCALIZAÇÃO|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "number_id": r"NUMBER ID\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "localizacao": r"LOCALIZAÇÃO\s*:\s*(.*?)(?=NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "site": r"SITE\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "ponto": r"PONTO\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|SEX|IDADE|DATAS|OBS|$)",
        "sex": r"SEX\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|IDADE|DATAS|OBS|$)",
        "idade": r"IDADE\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|DATAS|OBS|$)",
        "datas": r"DATAS\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|OBS|$)",
        "obs": r"OBS\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|$)",
    }

    fields = {}
    for key, pat in patterns.items():
        m = re.search(pat, text, re.IGNORECASE | re.DOTALL)
        if m:
            # Extract the value and clean up extra whitespace
            value = m.group(1).strip()
            # Replace multiple whitespace with single space
            value = re.sub(r"\s+", " ", value)
            fields[key] = value
        else:
            fields[key] = ""

    return fields


def parse_date(date_str: str) -> str:
    """Parse date string and return in YYYY-MM-DD format.

    Handles various formats like DD/MM/YYYY, DD.MM.YYYY, etc.
    Returns empty string if parsing fails.
    """
    if not date_str or date_str.strip() == "":
        return ""

    try:
        parsed = pd.to_datetime(date_str, errors="coerce")
        if pd.notna(parsed):
            return parsed.strftime("%Y-%m-%d")
    except Exception:
        pass

    return ""


def extract_images_from_slide(slide: Slide, slide_num: int, output_dir: Path) -> list[dict]:
    """Extract images and text metadata from a single slide.

    Returns a list of dicts, one per image found on the slide.
    Saves the full image and stores crop coordinates.
    """
    slide_text = ""
    slide_images = []

    for shape in slide.shapes:
        # Extract text
        if shape.has_text_frame:
            slide_text += shape.text + "\n"

        # Extract images
        if shape.shape_type == 13:  # picture
            try:
                image = shape.image
                ext = image.ext
                img_filename = f"slide_{slide_num:03d}_img_{len(slide_images)+1:02d}.{ext}"
                img_path = output_dir / img_filename

                # Save full image
                with open(img_path, "wb") as f:
                    f.write(image.blob)

                # Get image dimensions and crop values
                try:
                    img = Image.open(io.BytesIO(image.blob))
                    width, height = img.size

                    # Get crop values (as fractions, 0.0 to 1.0)
                    crop_left = shape.crop_left if hasattr(shape, "crop_left") else 0.0
                    crop_top = shape.crop_top if hasattr(shape, "crop_top") else 0.0
                    crop_right = shape.crop_right if hasattr(shape, "crop_right") else 0.0
                    crop_bottom = shape.crop_bottom if hasattr(shape, "crop_bottom") else 0.0

                    # Calculate pixel coordinates
                    left = int(width * crop_left)
                    top = int(height * crop_top)
                    right = int(width * (1 - crop_right))
                    bottom = int(height * (1 - crop_bottom))

                    if crop_left > 0 or crop_top > 0 or crop_right > 0 or crop_bottom > 0:
                        logging.debug(
                            "Slide %d image %d crop: left=%d, top=%d, right=%d, bottom=%d (%.1f%%, %.1f%%, %.1f%%, %.1f%%)",
                            slide_num,
                            len(slide_images) + 1,
                            left,
                            top,
                            right,
                            bottom,
                            crop_left * 100,
                            crop_top * 100,
                            crop_right * 100,
                            crop_bottom * 100,
                        )
                except Exception as e:
                    logging.warning("Failed to get crop info from slide %d: %s", slide_num, e)
                    left, top, right, bottom = pd.NA, pd.NA, pd.NA, pd.NA
                    width, height = pd.NA, pd.NA

                slide_images.append(
                    {
                        "filename": img_filename,
                        "width": width,
                        "height": height,
                        "crop_left": left,
                        "crop_top": top,
                        "crop_right": right,
                        "crop_bottom": bottom,
                    }
                )
            except Exception as e:
                logging.warning("Failed to extract image from slide %d: %s", slide_num, e)

    fields = extract_fields(slide_text)

    # We extract dates separately because sometimes it does not have a label
    dates = extract_dates(slide_text)

    # Create one row per image
    rows = []
    for img_data in slide_images:
        rows.append(
            {
                "fields": fields,
                "image_filename": img_data["filename"],
                "image_width": img_data["width"],
                "image_height": img_data["height"],
                "crop_left": img_data["crop_left"],
                "crop_top": img_data["crop_top"],
                "crop_right": img_data["crop_right"],
                "crop_bottom": img_data["crop_bottom"],
                "slide_num": slide_num,
                "dates": dates,
            }
        )

    return rows


def create_dataframe_row(idx: int, row_data: dict, output_dir: Path, input_pptx: Path) -> dict:
    """Convert extracted data into the standard labels format."""
    fields = row_data["fields"]
    img_filename = row_data["image_filename"]

    # Map PPTX fields to standard columns
    jaguar_id: str = fields.get("number_id", "") or fields.get("name", "") or fields.get("nome", "")
    sex: str = fields.get("sex", "")
    site_nr: str = fields.get("site") or fields.get("ponto")
    location: str = fields.get("localizacao", "")
    date_str: str = fields.get("datas", "")
    notes: str = fields.get("obs", "")

    # Parse Jaguar ID
    if jaguar_id:
        if "?" in jaguar_id:
            jaguar_id = jaguar_id.replace("?", "")

        jaguar_id = jaguar_id.strip().upper()
        # Normalize special characters: é -> E, ñ -> N, etc.
        jaguar_id = unicodedata.normalize("NFD", jaguar_id)
        jaguar_id = jaguar_id.encode("ascii", "ignore").decode("ascii")

        # if it ends with anything in brackets just remove that part and strip again
        bracket_regex = r"^(.*?)(\s*\(.*\))$"
        bracket_match = re.match(bracket_regex, jaguar_id)
        if bracket_match:
            jaguar_id = bracket_match.group(1).strip()
    else:
        jaguar_id = pd.NA

    # Parse site
    if site_nr:
        # remove leading zeros
        if site_nr.strip().startswith("0"):
            site_nr = site_nr.lstrip("0")
        site = site_nr if site_nr.upper().startswith("SITE") else f"SITE {site_nr}"
    else:
        site = pd.NA

    # Parse sex
    sexes = {
        "UNKNOWN": "U",
        "FEMALE JUVENILE": "F",
        "MALE JUVENILE": "M",
        "UNKNOWN JUVENILE": "U",
        "MALE": "M",
        "FEMALE": "F",
        "MACHO": "M",
        "FEMEA": "F",
        "M": "M",
        "F": "F",
        "U": "U",
    }

    for long, short in sexes.items():
        if sex.strip().upper().startswith(long):
            sex = short
            break

    if sex is None and pd.notna(jaguar_id):
        sex_regex = r"^C?N?(U|F|M)\d+.*$"
        sex_id_match = re.match(sex_regex, jaguar_id)
        sex = sex_id_match.group(1) if sex_id_match else pd.NA

    # Parse date
    date = parse_date(date_str)
    # only if it contains a single date elsewhere in the slide we use it
    if date == "" and len(row_data["dates"]) == 1:
        date = parse_date(row_data["dates"][0])

    # Sighting ID
    sighting_id = "PP" + str(idx + 1)

    # File metadata
    file_path = Path(img_filename)
    file_ext = file_path.suffix.lower()
    file_type = "IMAGE" if file_ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"] else "UNKNOWN"

    # Build the row
    return {
        "CAMERA TRAP SITE": site,
        "LATITUDE": pd.NA,
        "LONGITUDE": pd.NA,
        "CAMERA ID": pd.NA,
        "CAM": pd.NA,
        "JAGUAR ID": jaguar_id,
        "SEX": sex,
        "LOCATION": location.upper() if location else pd.NA,
        "CAMERA MODEL": pd.NA,
        "DATE": date,
        "TIME": pd.NA,
        "TEMP C": pd.NA,
        "Files Name": img_filename,
        "DATETIME": pd.NaT if not date else pd.to_datetime(date + " 00:00:00"),
        "ORIGINAL FILE NAME": img_filename,
        "SIGHTING ID": sighting_id,
        "RAW DATA PATH": input_pptx.as_posix(),
        "DATASET PATH": output_dir.as_posix(),
        "FILE PATH": file_path.as_posix(),
        "FILE TYPE": file_type,
        "FILE EXTENSION": file_ext,
        "FILE NAME": img_filename,
        "IMAGE WIDTH": row_data["image_width"],
        "IMAGE HEIGHT": row_data["image_height"],
        "CROP LEFT": row_data["crop_left"],
        "CROP TOP": row_data["crop_top"],
        "CROP RIGHT": row_data["crop_right"],
        "CROP BOTTOM": row_data["crop_bottom"],
        "ORIGINAL CAM": pd.NA,
        "ORIGINAL SITE": site,
        "NOTES": notes,
    }


def fill_missing_jaguar_ids(df: pd.DataFrame) -> pd.DataFrame:
    """Fill missing JAGUAR ID entries with sex-based IDs (F1, M1, U1, etc.).

    For each sex category (F, M, U), find the highest existing number and assign
    new sequential IDs to jaguars without names.
    """
    missing_id_mask = df["JAGUAR ID"].isna() | (df["JAGUAR ID"] == "") | (df["JAGUAR ID"] == pd.NA)
    num_missing = missing_id_mask.sum()

    if num_missing == 0:
        logging.info("No missing JAGUAR IDs to fill")
        return df

    logging.info("Filling %d missing JAGUAR IDs with sex-based identifiers", num_missing)

    # Find max existing numbers for each sex prefix
    max_nums = {"F": 0, "M": 0, "U": 0}

    id_regex = r"^(U|F|M)(\d+).*$"
    for jag_id in df.loc[~missing_id_mask, "JAGUAR ID"].astype(str):
        sex_id_match = re.match(id_regex, jag_id)

        if sex_id_match:
            sex = sex_id_match.group(1)
            id = sex_id_match.group(2)

            try:
                num = int(id)
                if num > max_nums[sex]:
                    max_nums[sex] = num
            except ValueError:
                continue

    # Assign new IDs based on sex
    for idx in df[missing_id_mask].index:
        sex = df.loc[idx, "SEX"]

        # Determine prefix based on sex
        if pd.isna(sex):
            prefix = "U"
        elif sex == "F":
            prefix = "F"
        elif sex == "M":
            prefix = "M"
        else:
            prefix = "U"

        # Increment counter and assign ID
        max_nums[prefix] += 1
        new_id = f"{prefix}{max_nums[prefix]}"
        df.loc[idx, "JAGUAR ID"] = new_id
        logging.debug("Assigned ID %s to row %d (sex=%s)", new_id, idx, sex)

    logging.info(
        "Assigned %d new JAGUAR IDs: F=%d, M=%d, U=%d",
        num_missing,
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if id.startswith("F")),
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if id.startswith("M")),
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if id.startswith("U")),
    )

    return df


def extract_from_pptx(input_pptx: Path, output_dir: Path) -> tuple[pd.DataFrame, dict]:
    """Extract all jaguar sighting data from PowerPoint file.

    Returns:
        - DataFrame with one row per image
        - Report dict with processing statistics
    """
    logging.info("Loading PowerPoint from %s", input_pptx)

    if not input_pptx.exists():
        raise FileNotFoundError(f"Input PowerPoint file not found: {input_pptx}")

    # Create output directory for images
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(input_pptx))

    all_rows = []
    num_slides_processed = 0
    num_images_extracted = 0
    idx = 0

    # Skip first slide (title slide)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue  # Skip title slide

        slide_rows = extract_images_from_slide(slide, i, output_dir)
        num_images_extracted += len(slide_rows)

        if slide_rows:
            num_slides_processed += 1
            for row_data in slide_rows:
                df_row = create_dataframe_row(idx, row_data, output_dir, input_pptx)
                all_rows.append(df_row)
                idx += 1

    df = pd.DataFrame(all_rows)

    # Fill missing JAGUAR IDs with sex-based IDs (F1, M1, U1, etc.)
    df = fill_missing_jaguar_ids(df)

    report = {
        "num_slides_total": len(prs.slides),
        "num_slides_with_images": num_slides_processed,
        "num_images_extracted": num_images_extracted,
        "num_rows_created": len(df),
        "unique_jaguar_ids": df["JAGUAR ID"].nunique() if len(df) > 0 else 0,
        "unique_sites": df["CAMERA TRAP SITE"].nunique() if len(df) > 0 else 0,
    }

    logging.info("Extracted %d images from %d slides", num_images_extracted, num_slides_processed)

    return df, report


def run(
    input_pptx: Path,
    output_csv: Path,
    output_dir: Path = None,
    generate_report: bool = False,
):
    """Main processing function."""
    if output_dir is None:
        output_dir = input_pptx.parent / "data"

    report = {}
    df, report["extraction"] = extract_from_pptx(input_pptx, output_dir)
    report["output"] = write_df(df, output_csv)

    if generate_report:
        report_path = output_csv.parent / "pptx_extraction_report.json"
        write_report(report, report_path)


def main():
    parser = argparse.ArgumentParser(description="Extract jaguar identification data from PowerPoint presentations.")
    parser.add_argument(
        "--input_pptx",
        type=str,
        required=True,
        help="Path to the input PowerPoint file.",
    )
    parser.add_argument(
        "--output_csv",
        type=str,
        required=True,
        help="Path to save the extracted labels CSV file.",
    )
    parser.add_argument(
        "--output_dir",
        type=str,
        default=None,
        help="Directory to save extracted images. Defaults to <input_pptx_parent>/data",
    )
    parser.add_argument(
        "--generate_report",
        action="store_true",
        help="Generate a processing report.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose logging.",
    )

    args = parser.parse_args()
    level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(level=level, format="%(levelname)s: %(message)s")

    output_dir = Path(args.output_dir) if args.output_dir else None

    run(
        Path(args.input_pptx),
        Path(args.output_csv),
        output_dir=output_dir,
        generate_report=args.generate_report,
    )


if __name__ == "__main__":
    raise SystemExit(main())
