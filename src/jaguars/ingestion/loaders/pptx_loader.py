"""PPTX ingestion loader module.

Usage (via CLI):
    # Default run command:
    python src/jaguars/ingestion/loaders/pptx_loader.py \
        --input-path "data/raw/17_11_2025/CAMERA TRAP ID GUIDE UPDATED by Oscar2025.pptx" \
        --verbose

    # With all options:
    python src/jaguars/ingestion/loaders/pptx_loader.py \
        --dataset "JID_Master_Dataset" \
        --input-path "data/raw/17_11_2025/CAMERA TRAP ID GUIDE UPDATED by Oscar2025.pptx" \
        --media-dir "data/intermediate/v1/pptx_media" \
        --output-csv "data/intermediate/v1/pptx_extracted_labels.csv" \
        --generate-report \
        --verbose
"""

import argparse
import io
import json
import logging
import re
import unicodedata
from pathlib import Path
from typing import Any, cast

import fiftyone as fo  # type: ignore
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.slide import Slide

from jaguars.common.config import DEFAULT_GROUP_SLICE, GROUP_FIELD_NAME, JID_MASTER_DATASET
from jaguars.common.fiftyone_utils import get_or_create_dataset
from jaguars.common.io_utils import ensure_dir
from jaguars.common.logging_utils import setup_logger

MODULE_NAME = "ingestion.loaders.pptx_loader"
logger = setup_logger(MODULE_NAME)


def validate_resources(input_path: Path) -> None:
    """Checks if inputs exist."""
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")


def write_summary(summary_data: dict[str, Any], summary_path: Path) -> None:
    """Generates summary report."""
    summary_path.parent.mkdir(parents=True, exist_ok=True)
    with open(summary_path, "w") as f:
        json.dump(summary_data, f, indent=4, default=str)


def run_processing(
    dataset_name: str,
    input_path: Path,
    media_dir: Path | None = None,
    output_csv: Path | None = None,
    generate_report: bool = False,
    detections_field: str = "pptx_detections",
    summary_location: str | None = None,
    dry_run: bool = False,
    verbose: bool = False,
) -> Any:
    """Core Logic for PPTX ingestion.

    Wrapper around ingest_pptx_slides to match the common module pattern.
    """
    logger = setup_logger(MODULE_NAME, level=logging.DEBUG if verbose else logging.INFO)
    logger.info("Starting processing for dataset: %s", dataset_name)

    validate_resources(input_path)

    if dry_run:
        logger.info("Dry run enabled - no changes made.")
        return None

    # Call the actual implementation
    dataset = ingest_pptx_slides(
        pptx_path=input_path,
        dataset_name=dataset_name,
        media_dir=media_dir,
        output_csv=output_csv,
        generate_report=generate_report,
        detections_field=detections_field,
    )

    if summary_location:
        summary_data = {
            "status": "success",
            "dataset_name": dataset_name,
            "input_path": str(input_path),
            "media_dir": str(media_dir),
            "processed": True,
        }
        write_summary(summary_data, Path(summary_location))

    logger.info("Processing completed successfully.")
    return dataset


def _pptx_media_dir(dataset: Any, dataset_name: str) -> Path:
    # Store derived media in a stable location associated with the FiftyOne dataset.
    # This avoids requiring a user-provided output directory.
    base = Path(getattr(fo.config, "default_dataset_dir", "."))
    return base / dataset_name / "media" / "pptx"


def _safe_int(v: Any) -> int | None:
    try:
        if pd.isna(v):
            return None
        return int(v)
    except Exception:
        return None


def _build_crop_detection(row: pd.Series) -> Any:
    width = _safe_int(row.get("IMAGE WIDTH"))
    height = _safe_int(row.get("IMAGE HEIGHT"))
    left = _safe_int(row.get("CROP LEFT"))
    top = _safe_int(row.get("CROP TOP"))
    right = _safe_int(row.get("CROP RIGHT"))
    bottom = _safe_int(row.get("CROP BOTTOM"))
    if not (width and height and left is not None and top is not None and right is not None and bottom is not None):
        return None
    if right <= left or bottom <= top:
        return None

    x = left / width
    y = top / height
    w = (right - left) / width
    h = (bottom - top) / height
    # Clamp to [0, 1]
    x = max(0.0, min(1.0, float(x)))
    y = max(0.0, min(1.0, float(y)))
    w = max(0.0, min(1.0 - x, float(w)))
    h = max(0.0, min(1.0 - y, float(h)))
    det = fo.Detection(label="pptx_crop", bounding_box=[x, y, w, h])
    return fo.Detections(detections=[det])


def ingest_pptx_slides(
    pptx_path: Path,
    dataset_name: str = JID_MASTER_DATASET,
    media_dir: Path | None = None,
    output_csv: Path | None = None,
    generate_report: bool = False,
    detections_field: str = "pptx_detections",
) -> Any:
    """Extract jaguar images + metadata from a PowerPoint and ingest into FiftyOne.

    Notes:
    - Extraction behavior matches legacy `pptx_extract.py`.
    - If `output_csv` is provided, a CSV in the legacy format is written.

    Args:
        pptx_path: Path to the input PowerPoint file.
        dataset_name: FiftyOne dataset to add samples into.
        media_dir: Optional directory to store extracted images for the dataset.
            If not provided, a default directory under FiftyOne's dataset dir is used.
        output_csv: Optional CSV output path (legacy format).
        generate_report: If True and output_csv is provided, write a JSON report next to it.
        detections_field: FiftyOne field name to store crop boxes as `fo.Detections`.
    """
    logger.info("Loading PowerPoint from %s", pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"Input PowerPoint file not found: {pptx_path}")

    dataset = get_or_create_dataset(dataset_name)

    # Set up group field if not already configured
    if dataset.group_field is None:
        dataset.add_group_field(GROUP_FIELD_NAME, default=DEFAULT_GROUP_SLICE)
        logger.info("Configured grouped dataset with 'image' and 'video' slices")

    if media_dir is None:
        media_dir = _pptx_media_dir(dataset, dataset_name)
    ensure_dir(media_dir)

    df, report, blobs = extract_from_pptx(pptx_path)

    # Keep legacy label schema consistent
    if len(df) > 0:
        df["DATASET PATH"] = media_dir.as_posix()

    # Persist extracted images into the dataset media directory
    for filename, blob in blobs.items():
        out_path = media_dir / filename
        if not out_path.exists():
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_bytes(blob)

    # Optionally write the legacy CSV output
    if output_csv is not None:
        output_csv.parent.mkdir(parents=True, exist_ok=True)
        df.to_csv(output_csv, index=False)
        logger.info("Writing extracted labels to %s", output_csv)
        if generate_report:
            report_path = output_csv.parent / "pptx_extraction_report.json"
            with open(report_path, "w") as f:
                json.dump(report, f, indent=4)
            logger.info("Writing processing report to %s", report_path)

    existing_uids = set()
    try:
        for v in dataset.values("pptx_uid"):
            if v:
                existing_uids.add(str(v))
    except Exception:
        existing_uids = set()

    # Ingest into FiftyOne (one sample per extracted image)
    new_samples: list[fo.Sample] = []
    for _, row in df.iterrows():
        img_filename = row["Files Name"]
        image_path = media_dir / str(img_filename)

        pptx_uid = f"{pptx_path.resolve().as_posix()}|{img_filename}"
        if pptx_uid in existing_uids:
            continue

        # Each image gets its own group ID in the 'image' slice
        import uuid

        group_id = str(uuid.uuid4())
        sample = fo.Sample(filepath=str(image_path), group=fo.Group().element(group_id))
        sample.group.name = DEFAULT_GROUP_SLICE

        sample["source_type"] = "pptx"
        sample["source"] = "pptx_camera_trap_guide_17_11_2025"
        sample["pptx_source"] = pptx_path.as_posix()
        sample["pptx_uid"] = pptx_uid
        m = re.match(r"^slide_(\d+)_img_\d+\.[A-Za-z0-9]+$", str(img_filename))
        sample["slide_num"] = int(m.group(1)) if m else None

        # Store the canonical fields used downstream
        if pd.notna(row.get("JAGUAR ID")):
            sample["jaguar_id"] = str(row["JAGUAR ID"])
            sample["ground_truth"] = fo.Classification(label=str(row["JAGUAR ID"]))
        if pd.notna(row.get("CAMERA TRAP SITE")):
            sample["site"] = str(row["CAMERA TRAP SITE"])
        if pd.notna(row.get("SEX")):
            sample["sex"] = str(row["SEX"])
        if pd.notna(row.get("LOCATION")):
            sample["location"] = str(row["LOCATION"])
        if pd.notna(row.get("DATE")):
            sample["date"] = str(row["DATE"])
        if pd.notna(row.get("NOTES")):
            sample["notes"] = str(row["NOTES"])

        # Preserve crop metadata
        for field in [
            "IMAGE WIDTH",
            "IMAGE HEIGHT",
            "CROP LEFT",
            "CROP TOP",
            "CROP RIGHT",
            "CROP BOTTOM",
        ]:
            if field in row and pd.notna(row[field]):
                sample[field.lower().replace(" ", "_")] = row[field]

        dets = _build_crop_detection(row)
        if dets is not None:
            sample[detections_field] = dets

        new_samples.append(sample)

    if new_samples:
        dataset.add_samples(new_samples)
        dataset.save()
        logger.info("Added %d samples from PPTX", len(new_samples))
    else:
        logger.info("No new samples added from PPTX")

    # Log ingestion summary
    summary = {
        "dataset_name": dataset_name,
        "pptx_source": str(pptx_path),
        "media_dir": str(media_dir),
        "slides_total": report.get("num_slides_total", 0),
        "slides_with_images": report.get("num_slides_with_images", 0),
        "images_extracted": report.get("num_images_extracted", 0),
        "new_samples_added": len(new_samples),
        "total_samples_in_dataset": len(dataset),
        "unique_jaguar_ids": report.get("unique_jaguar_ids", 0),
        "unique_sites": report.get("unique_sites", 0),
    }
    logger.info("PPTX Ingestion Summary: %s", json.dumps(summary, indent=2))

    return dataset


def extract_dates(text: str) -> list[str]:
    """Extract date information from slide text using regex patterns."""
    date_patterns = [
        r"(\d{1,2}[\/\.-]\d{1,2}[\/\.-]\d{2,4})",  # DD/MM/YYYY or DD-MM-YYYY or DD.MM.YYYY
        r"(\d{4}[\/\.-]\d{1,2}[\/\.-]\d{1,2})",  # YYYY/MM/DD or YYYY-MM-DD or YYYY.MM.DD
        r"(\d{1,2}\s+[A-Za-z]{3,9}\s+\d{2,4})",  # DD Month YYYY
        r"([A-Za-z]{3,9}\s+\d{1,2},\s+\d{4})",  # Month DD, YYYY
    ]

    dates_found = []
    for pattern in date_patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            dates_found.append(match.strip())

    return dates_found


def extract_fields(text: str) -> dict[str, str]:
    """Extract jaguar metadata fields from slide text using regex patterns."""
    patterns = {
        "nome": r"NOME\s*:\s*(.*?)(?=LOCALIZAÇÃO|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "name": r"NAME\s*:\s*(.*?)(?=LOCALIZAÇÃO|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "number_id": r"NUMBER ID\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "localizacao": r"LOCALIZAÇÃO\s*:\s*(.*?)(?=NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "site": r"SITE\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|PONTO|SEX|IDADE|DATAS|OBS|$)",
        "ponto": r"PONTO\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|SEX|IDADE|DATAS|OBS|$)",
        "sex": r"SEX\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|IDADE|DATAS|OBS|$)",
        "idade": r"IDADE\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|DATAS|OBS|$)",
        "datas": r"DATAS\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|OBS|$)",
        "obs": r"OBS\s*:\s*(.*?)(?=LOCALIZAÇÃO|NAME|NOME|NUMBER ID|SITE|PONTO|SEX|IDADE|DATAS|$)",
    }

    fields = {}
    for key, pat in patterns.items():
        m = re.search(pat, text, re.IGNORECASE | re.DOTALL)
        if m:
            value = m.group(1).strip()
            value = re.sub(r"\s+", " ", value)
            fields[key] = value
        else:
            fields[key] = ""

    return fields


def parse_date(date_str: str) -> str:
    """Parse date string and return in YYYY-MM-DD format.

    Handles various formats like DD/MM/YYYY, DD.MM.YYYY, etc.
    Returns empty string if parsing fails.
    """
    if not date_str or date_str.strip() == "":
        return ""

    try:
        parsed = pd.to_datetime(date_str, errors="coerce")
        if pd.notna(parsed):
            return str(parsed.strftime("%Y-%m-%d"))
    except Exception:
        pass

    return ""


def extract_images_from_slide(slide: Slide, slide_num: int) -> list[dict[str, Any]]:
    """Extract images and text metadata from a single slide.

    Returns a list of dicts, one per image found on the slide.
    Stores raw image bytes and crop coordinates.
    """
    slide_text = ""
    slide_images: list[dict[str, Any]] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            slide_text += str(getattr(shape, "text", "")) + "\n"

        if shape.shape_type == 13:  # picture
            try:
                image = cast(Picture, shape).image
                ext = image.ext
                img_filename = f"slide_{slide_num:03d}_img_{len(slide_images)+1:02d}.{ext}"
                blob = image.blob

                try:
                    img = Image.open(io.BytesIO(blob))
                    width, height = img.size

                    crop_left = shape.crop_left if hasattr(shape, "crop_left") else 0.0
                    crop_top = shape.crop_top if hasattr(shape, "crop_top") else 0.0
                    crop_right = shape.crop_right if hasattr(shape, "crop_right") else 0.0
                    crop_bottom = shape.crop_bottom if hasattr(shape, "crop_bottom") else 0.0

                    left = int(width * crop_left)
                    top = int(height * crop_top)
                    right = int(width * (1 - crop_right))
                    bottom = int(height * (1 - crop_bottom))

                    if crop_left > 0 or crop_top > 0 or crop_right > 0 or crop_bottom > 0:
                        logger.debug(
                            "Slide %d image %d crop: left=%d, top=%d, right=%d, bottom=%d (%.1f%%, %.1f%%, %.1f%%, %.1f%%)",
                            slide_num,
                            len(slide_images) + 1,
                            left,
                            top,
                            right,
                            bottom,
                            crop_left * 100,
                            crop_top * 100,
                            crop_right * 100,
                            crop_bottom * 100,
                        )
                except Exception as e:
                    logger.warning("Failed to get crop info from slide %d: %s", slide_num, e)
                    left, top, right, bottom = pd.NA, pd.NA, pd.NA, pd.NA
                    width, height = pd.NA, pd.NA

                slide_images.append(
                    {
                        "filename": img_filename,
                        "blob": blob,
                        "width": width,
                        "height": height,
                        "crop_left": left,
                        "crop_top": top,
                        "crop_right": right,
                        "crop_bottom": bottom,
                    }
                )
            except Exception as e:
                logger.warning("Failed to extract image from slide %d: %s", slide_num, e)

    fields = extract_fields(slide_text)
    dates = extract_dates(slide_text)

    rows = []
    for img_data in slide_images:
        rows.append(
            {
                "fields": fields,
                "image_filename": img_data["filename"],
                "blob": img_data["blob"],
                "image_width": img_data["width"],
                "image_height": img_data["height"],
                "crop_left": img_data["crop_left"],
                "crop_top": img_data["crop_top"],
                "crop_right": img_data["crop_right"],
                "crop_bottom": img_data["crop_bottom"],
                "slide_num": slide_num,
                "dates": dates,
            }
        )

    return rows


def create_dataframe_row(idx: int, row_data: dict[str, Any], output_dir: Path, input_pptx: Path) -> dict[str, Any]:
    """Convert extracted data into the standard labels format."""
    fields = row_data["fields"]
    img_filename = row_data["image_filename"]

    jaguar_id: str = fields.get("number_id", "") or fields.get("name", "") or fields.get("nome", "")
    sex: str = fields.get("sex", "")
    site_nr: str = fields.get("site") or fields.get("ponto")
    location: str = fields.get("localizacao", "")
    date_str: str = fields.get("datas", "")
    notes: str = fields.get("obs", "")

    if jaguar_id:
        if "?" in jaguar_id:
            jaguar_id = jaguar_id.replace("?", "")

        jaguar_id = jaguar_id.strip().upper()
        jaguar_id = unicodedata.normalize("NFD", jaguar_id)
        jaguar_id = jaguar_id.encode("ascii", "ignore").decode("ascii")

        bracket_regex = r"^(.*?)(\s*\(.*\))$"
        bracket_match = re.match(bracket_regex, jaguar_id)
        if bracket_match:
            jaguar_id = bracket_match.group(1).strip()
    else:
        jaguar_id = pd.NA

    if site_nr:
        if site_nr.strip().startswith("0"):
            site_nr = site_nr.lstrip("0")
        site = site_nr if site_nr.upper().startswith("SITE") else f"SITE {site_nr}"
    else:
        site = pd.NA

    sexes = {
        "UNKNOWN": "U",
        "FEMALE JUVENILE": "F",
        "MALE JUVENILE": "M",
        "UNKNOWN JUVENILE": "U",
        "MALE": "M",
        "FEMALE": "F",
        "MACHO": "M",
        "FEMEA": "F",
        "M": "M",
        "F": "F",
        "U": "U",
    }

    for long, short in sexes.items():
        if sex.strip().upper().startswith(long):
            sex = short
            break

    if sex is None and pd.notna(jaguar_id):
        sex_regex = r"^C?N?(U|F|M)\d+.*$"
        sex_id_match = re.match(sex_regex, jaguar_id)
        sex = sex_id_match.group(1) if sex_id_match else pd.NA

    date = parse_date(date_str)
    if date == "" and len(row_data["dates"]) == 1:
        date = parse_date(row_data["dates"][0])

    sighting_id = "PP" + str(idx + 1)

    file_path = Path(img_filename)
    file_ext = file_path.suffix.lower()
    file_type = "IMAGE" if file_ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"] else "UNKNOWN"

    return {
        "CAMERA TRAP SITE": site,
        "LATITUDE": pd.NA,
        "LONGITUDE": pd.NA,
        "CAMERA ID": pd.NA,
        "CAM": pd.NA,
        "JAGUAR ID": jaguar_id,
        "SEX": sex,
        "LOCATION": location.upper() if location else pd.NA,
        "CAMERA MODEL": pd.NA,
        "DATE": date,
        "TIME": pd.NA,
        "TEMP C": pd.NA,
        "Files Name": img_filename,
        "DATETIME": pd.NaT if not date else pd.to_datetime(date + " 00:00:00"),
        "ORIGINAL FILE NAME": img_filename,
        "SIGHTING ID": sighting_id,
        "RAW DATA PATH": input_pptx.as_posix(),
        "DATASET PATH": output_dir.as_posix(),
        "FILE PATH": file_path.as_posix(),
        "FILE TYPE": file_type,
        "FILE EXTENSION": file_ext,
        "FILE NAME": img_filename,
        "IMAGE WIDTH": row_data["image_width"],
        "IMAGE HEIGHT": row_data["image_height"],
        "CROP LEFT": row_data["crop_left"],
        "CROP TOP": row_data["crop_top"],
        "CROP RIGHT": row_data["crop_right"],
        "CROP BOTTOM": row_data["crop_bottom"],
        "ORIGINAL CAM": pd.NA,
        "ORIGINAL SITE": site,
        "NOTES": notes,
    }


def fill_missing_jaguar_ids(df: pd.DataFrame) -> pd.DataFrame:
    """Fill missing JAGUAR ID entries with sex-based IDs (F1, M1, U1, etc.)."""
    missing_id_mask = df["JAGUAR ID"].isna() | (df["JAGUAR ID"] == "") | (df["JAGUAR ID"] == pd.NA)
    num_missing = missing_id_mask.sum()

    if num_missing == 0:
        logger.info("No missing JAGUAR IDs to fill")
        return df

    logger.info("Filling %d missing JAGUAR IDs with sex-based identifiers", num_missing)

    max_nums = {"F": 0, "M": 0, "U": 0}

    id_regex = r"^(U|F|M)(\d+).*$"
    for jag_id in df.loc[~missing_id_mask, "JAGUAR ID"].astype(str):
        sex_id_match = re.match(id_regex, jag_id)

        if sex_id_match:
            sex = sex_id_match.group(1)
            id = sex_id_match.group(2)

            try:
                num = int(id)
                if num > max_nums[sex]:
                    max_nums[sex] = num
            except ValueError:
                continue

    for idx in df[missing_id_mask].index:
        sex = df.loc[idx, "SEX"]

        if pd.isna(sex):
            prefix = "U"
        elif sex == "F":
            prefix = "F"
        elif sex == "M":
            prefix = "M"
        else:
            prefix = "U"

        max_nums[prefix] += 1
        new_id = f"{prefix}{max_nums[prefix]}"
        df.loc[idx, "JAGUAR ID"] = new_id
        logger.debug("Assigned ID %s to row %d (sex=%s)", new_id, idx, sex)

    logger.info(
        "Assigned %d new JAGUAR IDs: F=%d, M=%d, U=%d",
        num_missing,
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if str(id).startswith("F")),
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if str(id).startswith("M")),
        sum(1 for id in df.loc[missing_id_mask, "JAGUAR ID"] if str(id).startswith("U")),
    )

    return df


def extract_from_pptx(input_pptx: Path) -> tuple[pd.DataFrame, dict[str, Any], dict[str, bytes]]:
    """Extract all jaguar sighting data from PowerPoint file."""
    logger.info("Loading PowerPoint from %s", input_pptx)

    if not input_pptx.exists():
        raise FileNotFoundError(f"Input PowerPoint file not found: {input_pptx}")

    prs = Presentation(str(input_pptx))

    all_rows = []
    blobs: dict[str, bytes] = {}
    num_slides_processed = 0
    num_images_extracted = 0
    idx = 0

    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue  # Skip title slide

        slide_rows = extract_images_from_slide(slide, i)
        num_images_extracted += len(slide_rows)

        if slide_rows:
            num_slides_processed += 1
            for row_data in slide_rows:
                # We fill DATASET PATH later from the ingestion function's media_dir
                # but keep legacy schema by passing a placeholder here.
                df_row = create_dataframe_row(idx, row_data, Path("."), input_pptx)
                all_rows.append(df_row)
                blobs[row_data["image_filename"]] = row_data["blob"]
                idx += 1

    df = pd.DataFrame(all_rows)

    if len(df) > 0:
        df = fill_missing_jaguar_ids(df)

    report = {
        "num_slides_total": len(prs.slides),
        "num_slides_with_images": num_slides_processed,
        "num_images_extracted": num_images_extracted,
        "num_rows_created": len(df),
        "unique_jaguar_ids": df["JAGUAR ID"].nunique() if len(df) > 0 else 0,
        "unique_sites": df["CAMERA TRAP SITE"].nunique() if len(df) > 0 else 0,
    }

    logger.info("Extracted %d images from %d slides", num_images_extracted, num_slides_processed)

    return df, report, blobs


def main() -> None:
    """CLI Entrypoint."""
    parser = argparse.ArgumentParser(description="PPTX Ingestion Loader")
    parser.add_argument("--dataset-name", default=JID_MASTER_DATASET, type=str, help="Name of the FiftyOne dataset to create or modify.")
    parser.add_argument("--input-path", required=True, type=Path, help="Path to input PPTX file.")
    parser.add_argument("--media-dir", type=Path, help="Directory to store extracted images.")
    parser.add_argument("--output-csv", type=Path, help="Path to write legacy CSV output.")
    parser.add_argument("--generate-report", action="store_true", help="Generate JSON report.")
    parser.add_argument("--detections-field", type=str, default="pptx_detections", help="FiftyOne field name for crop detections.")
    parser.add_argument("--summary-location", type=str, help="Location to save summary report.")
    parser.add_argument("--dry-run", action="store_true", help="If set, no changes will be made.")
    parser.add_argument("--verbose", action="store_true", help="Enable verbose logging output.")
    args = parser.parse_args()

    run_processing(
        dataset_name=args.dataset_name,
        input_path=args.input_path,
        media_dir=args.media_dir,
        output_csv=args.output_csv,
        generate_report=args.generate_report,
        detections_field=args.detections_field,
        summary_location=args.summary_location,
        dry_run=args.dry_run,
        verbose=args.verbose,
    )


if __name__ == "__main__":
    main()
