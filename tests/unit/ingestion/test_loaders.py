import shutil
import tempfile
import unittest
import uuid
from pathlib import Path

import fiftyone as fo
from PIL import Image
from pptx import Presentation

from jaguars.ingestion.loaders.csv_loader import ingest_csv_labels
from jaguars.ingestion.loaders.pptx_loader import ingest_pptx_slides


class TestLoaders(unittest.TestCase):
    def setUp(self) -> None:
        self.test_dir = Path(tempfile.mkdtemp())
        self.dataset_name = f"test_dataset_{uuid.uuid4().hex}"

    def tearDown(self) -> None:
        shutil.rmtree(self.test_dir, ignore_errors=True)
        for name in [self.dataset_name]:
            if name in fo.list_datasets():
                fo.delete_dataset(name)

    def test_ingest_pptx_adds_detections_and_dedupes(self) -> None:
        # Create tiny PNG
        img_path = self.test_dir / "img.png"
        Image.new("RGB", (100, 50), color=(255, 0, 0)).save(img_path)

        # Create PPTX with one title slide and one image slide
        pptx_path = self.test_dir / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])  # title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(str(img_path), 0, 0)
        # Apply crop
        pic.crop_left = 0.10
        pic.crop_top = 0.20
        pic.crop_right = 0.10
        pic.crop_bottom = 0.20
        prs.save(str(pptx_path))

        media_dir = self.test_dir / "pptx_media"

        ds = ingest_pptx_slides(
            pptx_path=pptx_path,
            dataset_name=self.dataset_name,
            media_dir=media_dir,
            detections_field="pptx_detections",
        )
        self.assertEqual(len(ds), 1)
        sample = ds.first()
        self.assertEqual(sample["source_type"], "pptx")
        self.assertIn("pptx_detections", sample)
        self.assertIsNotNone(sample["pptx_detections"])
        self.assertEqual(len(sample["pptx_detections"].detections), 1)

        # Re-run should not add duplicates
        ds2 = ingest_pptx_slides(
            pptx_path=pptx_path,
            dataset_name=self.dataset_name,
            media_dir=media_dir,
            detections_field="pptx_detections",
        )
        self.assertEqual(len(ds2), 1)

    def test_ingest_csv_ingests_images_and_videos_and_dedupes(self) -> None:
        # Build raw input_dir structure expected by legacy clean_labels
        input_dir = self.test_dir / "raw"
        video_dir = input_dir / "sites" / "SITE 1" / "CAM 1"
        image_dir = input_dir / "sites" / "SITE 1" / "ID PHOTOS"
        video_dir.mkdir(parents=True, exist_ok=True)
        image_dir.mkdir(parents=True, exist_ok=True)

        (video_dir / "test.mp4").touch()
        Image.new("RGB", (10, 10), color=(0, 255, 0)).save(image_dir / "test.jpg")

        # Minimal raw labels CSV (legacy will accept and derive file paths)
        labels_csv = input_dir / "labels.csv"
        labels_csv.write_text(
            "CAMERA TRAP SITE;LATITUDE;LONGITUDE;CAMERA ID;CAM;JAGUAR ID;LOCATION;CAMERA MODEL;DATE;TIME;TEMP C;Files Name;NOTES/ ERRORS\n"
            "SITE 1;;;;1;J1;LOC;;2025-01-01;;;test.mp4;\n"
            "SITE 1;;;;1;J2;LOC;;2025-01-01;;;test.jpg;\n",
            encoding="utf-8",
        )

        dataset = ingest_csv_labels(
            input_dir=input_dir,
            input_csv=Path("labels.csv"),
            dataset_name=self.dataset_name,
            auto_match_missing=False,
        )

        # Dataset should be grouped with image and video slices
        self.assertIsNotNone(dataset.group_field)

        # Check slices - each CSV row creates samples in the grouped dataset
        image_slice = dataset.select_group_slices("image")
        video_slice = dataset.select_group_slices("video")
        self.assertEqual(len(image_slice), 1)
        self.assertEqual(len(video_slice), 1)

        # Check metadata
        self.assertEqual(image_slice.first()["source_type"], "csv")
        self.assertEqual(video_slice.first()["source_type"], "csv")

        # Re-run should not add duplicates
        dataset2 = ingest_csv_labels(
            input_dir=input_dir,
            input_csv=Path("labels.csv"),
            dataset_name=self.dataset_name,
            auto_match_missing=False,
        )
        # Should still have same number of samples in each slice
        self.assertEqual(len(dataset2.select_group_slices("image")), 1)
        self.assertEqual(len(dataset2.select_group_slices("video")), 1)


if __name__ == "__main__":
    unittest.main()
